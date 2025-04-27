from fastapi import FastAPI, HTTPException, UploadFile, File, Query
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
import os
from dotenv import load_dotenv
import json
import re
import fitz  # PyMuPDF for PDFs
from pptx import Presentation  # for PPTX
from docx import Document  # for DOCX
import tempfile
import logging

# Load environment variables
load_dotenv()
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    raise Exception("GEMINI_API_KEY not found in environment variables!")
genai.configure(api_key=api_key)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# FastAPI app
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Restrict this in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

model = genai.GenerativeModel(model_name="gemini-1.5-flash")

class QuizRequest(BaseModel):
    topic: str
    num_questions: int
    difficulty: str

def parse_content_to_questions(content: str):
    # Cleaning and parsing content to structured questions
    content = re.sub(r"```.*?```", "", content, flags=re.DOTALL).strip()
    content = content.replace("```json", "").replace("```", "").strip()
    questions = []
    blocks = re.split(r"\n\s*\n", content)

    for block in blocks:
        if not block.strip():
            continue
        q_match = re.search(r"Question:\s*(.*)", block, re.IGNORECASE)
        options = { 'A': None, 'B': None, 'C': None, 'D': None }
        for opt in options.keys():
            opt_match = re.search(rf"{opt}\.\s*(.*)", block, re.IGNORECASE)
            if opt_match:
                options[opt] = opt_match.group(1).strip()
        correct_match = re.search(r"Correct Answer:\s*([A-D])", block, re.IGNORECASE)

        if q_match and all(options.values()) and correct_match:
            question = q_match.group(1).strip()
            correct_answer = correct_match.group(1).upper()
            questions.append({
                "question": question,
                "options": [
                    {"text": options['A'], "id": "A"},
                    {"text": options['B'], "id": "B"},
                    {"text": options['C'], "id": "C"},
                    {"text": options['D'], "id": "D"}
                ],
                "correctAnswer": correct_answer
            })
    return questions

@app.post("/generate_quiz")
async def generate_quiz(req: QuizRequest):
    if req.num_questions < 1 or req.num_questions > 10:
        raise HTTPException(status_code=400, detail="Number of questions must be between 1 and 10.")
    if req.difficulty not in ["easy", "medium", "hard"]:
        raise HTTPException(status_code=400, detail="Invalid difficulty level.")

    prompt = (
        f"Generate {req.num_questions} {req.difficulty} level multiple-choice quiz questions about {req.topic}.\n"
        "Each question should have 4 options (A, B, C, D) with exactly one correct answer.\n"
        "Use this format:\n\n"
        "Question: <your question>\n"
        "A. <option A>\n"
        "B. <option B>\n"
        "C. <option C>\n"
        "D. <option D>\n"
        "Correct Answer: <letter>\n\n"
        "Only respond in the format above."
    )

    try:
        response = model.generate_content(
            contents=[{"role": "user", "parts": [{"text": prompt}]}]
        )
        content = response.text
        questions = parse_content_to_questions(content)
        return {"questions": questions}
    except Exception as e:
        logger.error(f"Error generating quiz: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate quiz: {str(e)}")

def extract_text_from_file(file: UploadFile):
    if file.filename.endswith('.pdf'):
        with fitz.open(stream=file.file.read(), filetype="pdf") as doc:
            text = "\n".join(page.get_text() for page in doc)
    elif file.filename.endswith('.pptx'):
        prs = Presentation(file.file)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    elif file.filename.endswith('.docx'):
        doc = Document(file.file)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type. Upload PDF, PPTX, or DOCX.")
    return text

@app.post("/upload_and_generate_quiz")
async def upload_and_generate_quiz(
    file: UploadFile = File(...),
    num_questions: int = Query(5, ge=1, le=10),
    difficulty: str = Query("medium", regex="^(easy|medium|hard)$")
):
    try:
        text = extract_text_from_file(file)
        prompt = (
            f"Generate {num_questions} {difficulty} level multiple-choice quiz questions based on the following content:\n\n"
            f"{text}\n\n"
            "Each question should have 4 options (A, B, C, D) and one correct answer.\n"
            "Format:\n"
            "Question: <your question>\n"
            "A. <option A>\n"
            "B. <option B>\n"
            "C. <option C>\n"
            "D. <option D>\n"
            "Correct Answer: <letter>\n\n"
            "Respond only with questions and options."
        )
        response = model.generate_content(
            contents=[{"role": "user", "parts": [{"text": prompt}]}]
        )
        content = response.text
        questions = parse_content_to_questions(content)

        # Save to temporary file
        temp = tempfile.NamedTemporaryFile(delete=False, suffix=".json")
        with open(temp.name, 'w') as f:
            json.dump({"questions": questions}, f, indent=4)

        return {"questions": questions, "download_url": f"/download_quiz?file_path={temp.name}"}
    except Exception as e:
        logger.error(f"Error uploading file or generating quiz: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate quiz: {str(e)}")

@app.get("/download_quiz")
async def download_quiz(file_path: str):
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found.")
    return FileResponse(path=file_path, filename="quiz.json", media_type='application/json')
